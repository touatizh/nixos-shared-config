from pptx import Presentation
from pptx.util import Pt, Inches
from docx import Document
import sys
import os
import re

def add_numbered_paragraph(slide, text: str, top: float = 2.0) -> None:
    """
    Add a numbered paragraph to a slide as a text box.
    
    Args:
        slide: PowerPoint slide object.
        text (str): Text content to add.
        top (float): Distance from the top of the slide (in inches).
    """
    # Add a text box to the slide
    left = Inches(1.0)
    width = Inches(8.0)
    height = Inches(1.0)
    textbox = slide.shapes.add_textbox(left, Inches(top), width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(18)

def docx_to_pptx_by_heading1(docx_file: str, pptx_file: str) -> None:
    """
    Convert a Word (.docx) document into a PowerPoint (.pptx) presentation.
    Each Heading 1 starts a new slide. Numbered lists are preserved as separate paragraphs.

    Args:
        docx_file (str): Path to the input Word document.
        pptx_file (str): Path to the output PowerPoint file.

    Returns:
        None
    """
    # Load the Word document
    doc = Document(docx_file)
    
    # Create a PowerPoint presentation
    presentation = Presentation()

    # Initialize variables
    slide = None
    current_top = 2.0  # Tracks the vertical position for text placement

    # Process paragraphs in the Word document
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue  # Skip empty lines

        # Create a new slide for Heading 1
        if paragraph.style.name == "Heading 2":
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title_placeholder = slide.shapes.title
            title_placeholder.text = text  # Set slide title
            current_top = 2.0  # Reset vertical position for new slide
        else:
            if slide:
                # Process numbered lists and regular text
                if re.match(r"^\d+\.", text):  # Numbered list
                    add_numbered_paragraph(slide, text, current_top)
                else:  # Regular bullet point or text
                    add_numbered_paragraph(slide, f"• {text}", current_top)
                current_top += 0.5  # Increment vertical position for next line

    # Save the PowerPoint presentation
    presentation.save(pptx_file)
    print(f"Presentation successfully created: {pptx_file}")

if __name__ == "__main__":
    # Check for arguments
    if len(sys.argv) != 3:
        print("Usage: python generate_ppt_from_docx.py <input_docx_file> <output_pptx_file>")
        sys.exit(1)

    input_docx = sys.argv[1]
    output_pptx = sys.argv[2]

    if not os.path.isfile(input_docx):
        print(f"Error: Input file '{input_docx}' does not exist.")
        sys.exit(1)

    # Generate the PowerPoint file
    docx_to_pptx_by_heading1(input_docx, output_pptx)
